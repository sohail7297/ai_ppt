from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ---------- THEMES ----------

THEMES = {
    "Black Gold": {
        "title": RGBColor(255, 215, 0),
        "text": RGBColor(255, 255, 255),
        "bg": RGBColor(15, 15, 15)
    },
    "Corporate": {
        "title": RGBColor(21, 101, 192),
        "text": RGBColor(0, 0, 0),
        "bg": RGBColor(255, 255, 255)
    },
    "Startup Pitch": {
        "title": RGBColor(142, 36, 170),
        "text": RGBColor(0, 0, 0),
        "bg": RGBColor(255, 255, 255)
    },
    "Academic": {
        "title": RGBColor(46, 125, 50),
        "text": RGBColor(0, 0, 0),
        "bg": RGBColor(255, 255, 255)
    }
}


# ---------- HELPERS ----------

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_speaker_notes(slide, notes):
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except:
        pass


# ---------- MAIN GENERATOR ----------

def create_presentation(
    topic,
    slides_data,
    theme_name="Black Gold",
    output_dir="generated/ppt"
):

    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()

    theme = THEMES.get(theme_name, THEMES["Black Gold"])

    # =================================================
    # TITLE SLIDE
    # =================================================

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    set_slide_background(slide, theme["bg"])

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = topic

    title.text_frame.paragraphs[0].font.size = Pt(30)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = theme["title"]

    subtitle.text = "Generated using AI Presentation Studio"

    # =================================================
    # CONTENT SLIDES
    # =================================================

    for slide_info in slides_data:

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        set_slide_background(slide, theme["bg"])

        title = slide.shapes.title
        title.text = slide_info["title"]

        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = theme["title"]

        content = slide.placeholders[1]

        tf = content.text_frame
        tf.clear()

        first = True

        for point in slide_info["points"]:
            if first:
                tf.paragraphs[0].text = point
                first = False
            else:
                p = tf.add_paragraph()
                p.text = point
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = theme["text"]

        add_speaker_notes(
            slide,
            slide_info.get("speaker_notes", "")
        )

    # =================================================
    # THANK YOU SLIDE
    # =================================================

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    set_slide_background(slide, theme["bg"])

    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(1)
    )

    tf = textbox.text_frame
    tf.text = "Thank You"

    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = theme["title"]

    file_path = os.path.join(
        output_dir,
        f"{topic.replace(' ','_')}.pptx"
    )

    prs.save(file_path)

    return file_path